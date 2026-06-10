"""
多業務索引 + 凌晨同步腳本（GitHub Actions 每天執行）

每個啟用的業務（businesses.json）會：
  1. 同步：把「待審核」分頁中 狀態=已審查 且屬於該業務的列，搬到該業務正式分頁，
     並從待審核移除。
  2. 索引文件：掃描該業務的 Drive 資料夾 → 本地 fastembed → 寫入該業務 namespace。
     （掃描檔 PDF 會嘗試 OCR）
  3. 索引 FAQ：把該業務正式分頁的問答也建索引（source_type=faq），納入查詢。

需要的環境變數：
  GOOGLE_SERVICE_ACCOUNT_JSON, PINECONE_API_KEY, PINECONE_INDEX_NAME(選填)
"""

import io
import os
import re
import json
import time
import hashlib
import urllib.parse
import urllib.request
from datetime import datetime, timezone, timedelta

from google.oauth2 import service_account
import google.auth.transport.requests
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from pinecone import Pinecone, ServerlessSpec
from fastembed import TextEmbedding

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

# ── 設定 ─────────────────────────────────────────────────────────────────────
PINECONE_API_KEY = os.environ["PINECONE_API_KEY"]
INDEX_NAME       = os.environ.get("PINECONE_INDEX_NAME", "ah-biao-bot")
SA_INFO          = json.loads(os.environ["GOOGLE_SERVICE_ACCOUNT_JSON"])
EMBED_MODEL_NAME = "BAAI/bge-small-zh-v1.5"
EMBED_DIM        = 512
CHUNK_SIZE       = 800
CHUNK_OVERLAP    = 150
EMBED_BATCH      = 256
MAX_FILE_MB      = 20
OCR_MAX_PAGES    = 15        # 每個掃描 PDF 最多 OCR 頁數
OCR_MIN_TEXT     = 80        # 抽出文字少於此字數 → 視為掃描檔，嘗試 OCR
GEMINI_API_KEY   = os.environ.get("GEMINI_API_KEY", "")
OUTLINE_MIN_CHARS = 2500     # 文件還原後字數 ≥ 此值 → 視為教學/大型檔，生成大綱索引
GEMINI_MODELS    = ["gemini-2.5-flash", "gemini-2.0-flash", "gemini-flash-latest"]

with open(os.path.join(ROOT, "businesses.json"), "r", encoding="utf-8") as f:
    CONFIG = json.load(f)
FAQ_SHEET_ID = CONFIG["faq_sheet_id"]
REVIEW_TAB   = CONFIG.get("review_tab", "待審核")
BUSINESSES   = [b for b in CONFIG["businesses"] if b.get("enabled")]

SKIP_MIMES = {"image/", "video/", "audio/", "application/zip", "application/x-zip",
              "application/x-rar", "application/octet-stream", "application/vnd.ms-powerpoint"}
def should_skip(m): return any(m.startswith(x) for x in SKIP_MIMES)

def now_tw_date(): return datetime.now(timezone(timedelta(hours=8))).strftime("%Y-%m-%d")

# ── Google 憑證 ───────────────────────────────────────────────────────────────
def drive_service():
    creds = service_account.Credentials.from_service_account_info(
        SA_INFO, scopes=["https://www.googleapis.com/auth/drive.readonly"])
    return build("drive", "v3", credentials=creds)

def sheets_token():
    creds = service_account.Credentials.from_service_account_info(
        SA_INFO, scopes=["https://www.googleapis.com/auth/spreadsheets"])
    creds.refresh(google.auth.transport.requests.Request())
    return creds.token

# ── Sheets REST helpers ──────────────────────────────────────────────────────
def _sheet_req(method, path, body=None):
    token = sheets_token()
    url = f"https://sheets.googleapis.com/v4/spreadsheets/{path}"
    data = json.dumps(body).encode() if body is not None else None
    req = urllib.request.Request(url, data=data, method=method,
          headers={"Authorization": f"Bearer {token}", "Content-Type": "application/json"})
    with urllib.request.urlopen(req, timeout=30) as r:
        return json.loads(r.read())

def sheet_read(tab):
    rng = urllib.parse.quote(f"{tab}!A1:J5000")
    resp = _sheet_req("GET", f"{FAQ_SHEET_ID}/values/{rng}")
    return resp.get("values", [])

def sheet_append(tab, rows):
    if not rows:
        return
    rng = urllib.parse.quote(f"{tab}!A1")
    _sheet_req("POST",
        f"{FAQ_SHEET_ID}/values/{rng}:append?valueInputOption=USER_ENTERED&insertDataOption=INSERT_ROWS",
        {"values": rows})

def sheet_tab_id(tab):
    meta = _sheet_req("GET", f"{FAQ_SHEET_ID}?fields=sheets.properties")
    for s in meta.get("sheets", []):
        if s["properties"]["title"] == tab:
            return s["properties"]["sheetId"]
    return None

def sheet_delete_rows(tab, row_indices_zero_based):
    """刪除指定的列（0-based，含表頭列為第0列）。由大到小刪避免位移。"""
    if not row_indices_zero_based:
        return
    sid = sheet_tab_id(tab)
    if sid is None:
        return
    reqs = []
    for idx in sorted(row_indices_zero_based, reverse=True):
        reqs.append({"deleteDimension": {"range": {
            "sheetId": sid, "dimension": "ROWS",
            "startIndex": idx, "endIndex": idx + 1}}})
    _sheet_req("POST", f"{FAQ_SHEET_ID}:batchUpdate", {"requests": reqs})

# 索引問題清單（無法索引的檔案，會寫到試算表「索引問題」分頁供使用者處理）
_PROBLEMS = []
def _report_problem(name, reason, size_mb=0):
    _PROBLEMS.append({"name": name, "reason": reason, "size": round(size_mb, 1)})

# ── Drive 掃描 / 取文字 ───────────────────────────────────────────────────────
def list_files(service, folder_id, depth=0):
    files, page = [], None
    while True:
        resp = service.files().list(
            q=f"'{folder_id}' in parents and trashed=false",
            fields="nextPageToken, files(id,name,mimeType,modifiedTime,size)",
            pageToken=page, pageSize=100).execute()
        for f in resp.get("files", []):
            if f["mimeType"] == "application/vnd.google-apps.folder":
                files.extend(list_files(service, f["id"], depth + 1))
            elif not should_skip(f["mimeType"]):
                size_mb = int(f.get("size", 0)) / 1024 / 1024
                if size_mb <= MAX_FILE_MB:
                    files.append(f)
                else:
                    _report_problem(f["name"], f"檔案過大（{size_mb:.0f}MB，超過{MAX_FILE_MB}MB上限）→ 建議改放純文字版或拆檔", size_mb)
        page = resp.get("nextPageToken")
        if not page:
            break
    return files

def _download(service, fid, export_mime=None, expected_size=0):
    # 一次取回完整內容（避免分塊下載偶發截斷）
    if export_mime:
        data = service.files().export_media(fileId=fid, mimeType=export_mime).execute()
    else:
        data = service.files().get_media(fileId=fid).execute()
    if isinstance(data, str):
        data = data.encode("utf-8")
    # 防靜默截斷：實際下載量明顯小於檔案大小 → 視為下載失敗（觸發重試）
    if expected_size and len(data) < int(expected_size) * 0.98:
        raise IOError(f"下載不完整：{len(data)}/{expected_size} bytes")
    return data

def _pdf_text(raw, name):
    try:
        from pypdf import PdfReader
        r = PdfReader(io.BytesIO(raw))
        return "\n".join(p.extract_text() or "" for p in r.pages)
    except Exception as e:
        print(f"  ⚠ PDF 解析失敗 {name}: {e}")
        return ""

def _ocr_pdf(raw, name):
    """掃描 PDF 用 OCR 取文字（需 tesseract + poppler）。"""
    try:
        from pdf2image import convert_from_bytes
        import pytesseract
        images = convert_from_bytes(raw, dpi=200, first_page=1, last_page=OCR_MAX_PAGES)
        out = []
        for im in images:
            out.append(pytesseract.image_to_string(im, lang="chi_tra+eng"))
        text = "\n".join(out)
        if text.strip():
            print(f"  🔍 OCR 取得 {len(text)} 字：{name}")
        return text
    except Exception as e:
        print(f"  ⚠ OCR 失敗 {name}: {e}")
        return ""

def _docx_paras(doc):
    """最保險的純段落讀法（與舊版相同，當作後備）。"""
    return "\n".join(p.text for p in doc.paragraphs)

def _docx_text(raw, name):
    """讀 .docx：段落 + 表格（依文件順序），避免漏掉放在表格內的法條/契約條款。
    表格走訪若因 python-docx 版本等問題出錯，會自動退回純段落讀法，絕不回傳空字串。"""
    try:
        from docx import Document
        doc = Document(io.BytesIO(raw))
    except Exception as e:
        print(f"  ⚠ DOCX 開檔失敗 {name}: {e}")
        return ""

    # 先嘗試「段落+表格依序」的完整抽取
    try:
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        def cell_text(cell):
            return " ".join(p.text for p in cell.paragraphs if p.text.strip())

        def table_text(tbl):
            lines = []
            for row in tbl.rows:
                cells = [cell_text(c).strip() for c in row.cells]
                line = " | ".join(c for c in cells if c)
                if line:
                    lines.append(line)
            return "\n".join(lines)

        parts = []
        for child in doc.element.body.iterchildren():
            tag = child.tag.split("}")[-1]
            if tag == "p":
                t = Paragraph(child, doc).text
                if t.strip():
                    parts.append(t)
            elif tag == "tbl":
                t = table_text(Table(child, doc))
                if t.strip():
                    parts.append(t)
        full = "\n".join(parts)
        paras = _docx_paras(doc)
        # 取內容較多者，避免任何情況下比舊版還少
        text = full if len(full.strip()) >= len(paras.strip()) else paras
        if text.strip():
            return text
        return paras
    except Exception as e:
        # 表格走訪失敗 → 退回純段落（等同舊版行為，不退化為空）
        print(f"  ⚠ DOCX 表格抽取失敗，改用純段落 {name}: {e}")
        try:
            return _docx_paras(doc)
        except Exception as e2:
            print(f"  ⚠ DOCX 純段落也失敗 {name}: {e2}")
            return ""

def _doc_text(raw, name):
    """舊版 .doc：python-docx 讀不了，用 LibreOffice 轉純文字（中文友善）。"""
    import subprocess, tempfile
    try:
        with tempfile.TemporaryDirectory() as td:
            src = os.path.join(td, "f.doc")
            with open(src, "wb") as fh:
                fh.write(raw)
            env = dict(os.environ, HOME=td)
            subprocess.run(["libreoffice", "--headless", "--convert-to", "txt:Text",
                            "--outdir", td, src],
                           timeout=150, capture_output=True, env=env)
            out = os.path.join(td, "f.txt")
            if os.path.exists(out):
                with open(out, encoding="utf-8", errors="ignore") as fh:
                    return fh.read()
    except Exception as e:
        print(f"  ⚠ DOC 轉換失敗 {name}: {e}")
    return ""

def _xlsx_text(raw, name):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        rows = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                line = "\t".join(str(c) if c is not None else "" for c in row)
                if line.strip():
                    rows.append(line)
        return "\n".join(rows)
    except Exception as e:
        print(f"  ⚠ XLSX 失敗 {name}: {e}")
        return ""

def _pptx_text(raw, name):
    """讀 .pptx：每張投影片的文字方塊與表格（含備忘稿）。"""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw))
        out = []

        def shape_text(shp):
            parts = []
            if shp.has_text_frame:
                for p in shp.text_frame.paragraphs:
                    t = "".join(r.text for r in p.runs)
                    if t.strip():
                        parts.append(t)
            if shp.has_table:
                for row in shp.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    line = " | ".join(c for c in cells if c)
                    if line:
                        parts.append(line)
            # 群組圖形遞迴
            if shp.shape_type == 6:  # GROUP
                for s in shp.shapes:
                    parts.extend(shape_text(s))
            return parts

        for i, slide in enumerate(prs.slides, 1):
            sp = []
            for shp in slide.shapes:
                try:
                    sp.extend(shape_text(shp))
                except Exception:
                    pass
            if slide.has_notes_slide:
                nt = slide.notes_slide.notes_text_frame.text
                if nt and nt.strip():
                    sp.append("（備忘稿）" + nt.strip())
            if sp:
                out.append(f"投影片 {i}：\n" + "\n".join(sp))
        return "\n\n".join(out)
    except Exception as e:
        print(f"  ⚠ PPTX 失敗 {name}: {e}")
        return ""

XLSX_MIMES = {"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
              "application/vnd.ms-excel"}
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

def extract_text(service, f):
    mime, fid, name = f["mimeType"], f["id"], f["name"]
    sz = f.get("size", 0)   # 二進位檔的大小，用來驗證下載完整性
    try:
        if mime == "application/vnd.google-apps.document":
            return _download(service, fid, "text/plain").decode("utf-8", "ignore")
        if mime == "application/vnd.google-apps.spreadsheet":
            return _download(service, fid, "text/csv").decode("utf-8", "ignore")
        if mime == "application/vnd.google-apps.presentation":
            return _download(service, fid, "text/plain").decode("utf-8", "ignore")
        if mime == "application/pdf":
            raw = _download(service, fid, expected_size=sz)
            text = _pdf_text(raw, name)
            if len((text or "").strip()) < OCR_MIN_TEXT:   # 疑似掃描檔 → OCR
                text = _ocr_pdf(raw, name) or text
            return text
        if mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return _docx_text(_download(service, fid, expected_size=sz), name)
        if mime == "application/msword":          # 舊版 .doc
            return _doc_text(_download(service, fid, expected_size=sz), name)
        if mime == PPTX_MIME:
            return _pptx_text(_download(service, fid, expected_size=sz), name)
        if mime in XLSX_MIMES:
            return _xlsx_text(_download(service, fid, expected_size=sz), name)
        if mime.startswith("text/"):
            return _download(service, fid, expected_size=sz).decode("utf-8", "ignore")
    except Exception as e:
        print(f"  ⚠ 讀取失敗（暫時性，將重試）{name}: {e}")
        return None          # 例外＝暫時性錯誤 → 回 None 讓上層重試
    return ""                # 走到這＝不支援的型態（視為真空白）

def chunk_text(text, source, file_id, modified, namespace):
    text = (text or "").strip()
    if not text:
        return []
    out, start, idx = [], 0, 0
    prefix = f"{namespace}:" if namespace else ""   # 預設 namespace 用舊 ID 規則 → 覆蓋而非重複
    while start < len(text):
        end = min(start + CHUNK_SIZE, len(text))
        out.append({
            "id": hashlib.md5(f"{prefix}{file_id}:{idx}".encode()).hexdigest(),
            "text": text[start:end], "source": source,
            "file_id": file_id, "modified": modified, "chunk_idx": idx,
        })
        start += CHUNK_SIZE - CHUNK_OVERLAP
        idx += 1
    return out

# ── Embedding ────────────────────────────────────────────────────────────────
_model = None
def get_model():
    global _model
    if _model is None:
        print(f"  載入本地 embedding 模型 {EMBED_MODEL_NAME} ...")
        _model = TextEmbedding(model_name=EMBED_MODEL_NAME)
    return _model

def embed(texts):
    return [e.tolist() for e in get_model().embed(texts, batch_size=EMBED_BATCH)]

# ── Pinecone ─────────────────────────────────────────────────────────────────
def get_index(pc):
    if INDEX_NAME not in [i.name for i in pc.list_indexes()]:
        print(f"  建立 index {INDEX_NAME} (dim={EMBED_DIM})")
        pc.create_index(name=INDEX_NAME, dimension=EMBED_DIM, metric="cosine",
                        spec=ServerlessSpec(cloud="aws", region="us-east-1"))
        time.sleep(10)
    return pc.Index(INDEX_NAME)

def upsert(index, namespace, items, embeddings, source_type):
    vectors = []
    for it, emb in zip(items, embeddings):
        md = {
            "source": it["source"], "text": it["text"][:900],
            "source_type": source_type,
            "file_id": it.get("file_id", ""), "chunk_idx": it.get("chunk_idx", 0),
        }
        md.update(it.get("meta_extra", {}))   # FAQ 額外帶 faq_answer / faq_question
        vectors.append({"id": it["id"], "values": emb, "metadata": md})
    for i in range(0, len(vectors), 100):
        index.upsert(vectors=vectors[i:i+100], namespace=namespace)

# ── Checkpoint ───────────────────────────────────────────────────────────────
CKPT = os.path.join(ROOT, "data", "checkpoint.json")
def load_ckpt():
    if os.path.exists(CKPT):
        with open(CKPT) as f: return json.load(f)
    return {}
def save_ckpt(d):
    os.makedirs(os.path.dirname(CKPT), exist_ok=True)
    with open(CKPT, "w") as f: json.dump(d, f)

# ── 步驟一：同步已審查 → 正式分頁（掃描 待審核 + 其他素材）──────────────────────
PROMOTE_SOURCE_TABS = [REVIEW_TAB, "其他素材"]

def promote_reviewed(business):
    """從來源分頁把『業務類別=本業務 且 狀態=已審查』的列，搬到本業務正式分頁。"""
    for src_tab in PROMOTE_SOURCE_TABS:
        if src_tab == business["sheet_tab"]:
            continue
        try:
            rows = sheet_read(src_tab)
        except Exception:
            continue
        if not rows:
            continue
        move, del_idx = [], []
        for i in range(1, len(rows)):
            r = (rows[i] + [""] * 10)[:10]
            biz, status = r[1].strip(), r[8].strip()
            if biz == business["name"] and status == "已審查":
                today = now_tw_date()
                r[6] = r[6] or today
                r[7] = today
                move.append(r)
                del_idx.append(i)
        if move:
            sheet_append(business["sheet_tab"], move)
            sheet_delete_rows(src_tab, del_idx)
            print(f"  ➡ 從「{src_tab}」搬 {len(move)} 筆到「{business['sheet_tab']}」")

# ── 步驟二：索引 Drive 資料夾 ────────────────────────────────────────────────
def index_drive(index, business, drive, done):
    ns = business["namespace"]
    fid_folder = business.get("drive_folder_id", "")
    if not fid_folder:
        return 0
    files = list_files(drive, fid_folder)
    print(f"  Drive 找到 {len(files)} 個檔案")
    total = 0
    for i, f in enumerate(files, 1):
        fid, fmod = f["id"], f.get("modifiedTime", "")
        prev = done.get(fid, "")
        if prev == fmod:                       # 已成功索引（有內容）
            continue
        # 暫時性錯誤的重試次數：值格式 err<N>:<modifiedTime>
        sz_mb = int(f.get("size", 0)) / 1024 / 1024
        em = re.match(r"^err(\d+):(.*)$", prev) if prev else None
        retries = int(em.group(1)) if (em and em.group(2) == fmod) else 0
        if retries >= 5:                       # 暫時性錯誤最多重試 5 次才放棄 → 列入問題清單
            _report_problem(f["name"], "讀取一直失敗（已重試5次）→ 建議檢查檔案或重新上傳", sz_mb)
            continue
        text = extract_text(drive, f)
        if text is None:                       # 暫時性錯誤（下載/解析爆錯）→ 記次數，下次再試
            done[fid] = f"err{retries+1}:{fmod}"
            if retries + 1 >= 5:               # 這次達上限 → 列入問題清單
                _report_problem(f["name"], "讀取一直失敗（已重試5次）→ 建議檢查檔案或重新上傳", sz_mb)
            print(f"  [{i}/{len(files)}] {f['name']} → 讀取錯誤（重試 {retries+1}/5）")
            save_ckpt(done)
            continue
        chunks = chunk_text(text, f["name"], fid, fmod, ns)
        if chunks:
            for j in range(0, len(chunks), EMBED_BATCH):
                bc = chunks[j:j+EMBED_BATCH]
                upsert(index, ns, bc, embed([c["text"] for c in bc]), "doc")
            total += len(chunks)
            print(f"  [{i}/{len(files)}] {f['name']} → {len(chunks)} 片段")
        else:                                  # 抽不到任何文字 → 列入問題清單
            _report_problem(f["name"], "取不到文字（可能是圖片型/空白檔，連OCR都讀不到）→ 建議改放可選取文字的版本", sz_mb)
            print(f"  [{i}/{len(files)}] {f['name']} → 真的沒文字（標記完成）")
        done[fid] = fmod                       # 有內容或確定空白 → 標記完成
        save_ckpt(done)
    return total

# ── 把「索引問題清單」寫到試算表，供使用者處理 ────────────────────────────────
def write_problem_report():
    tab = "索引問題"
    header = ["檔名", "問題", "大小(MB)", "偵測時間"]
    try:
        meta = _sheet_req("GET", f"{FAQ_SHEET_ID}?fields=sheets.properties.title")
        tabs = [s["properties"]["title"] for s in meta.get("sheets", [])]
        if tab not in tabs:
            _sheet_req("POST", f"{FAQ_SHEET_ID}:batchUpdate",
                       {"requests": [{"addSheet": {"properties": {"title": tab,
                        "gridProperties": {"frozenRowCount": 1}}}}]})
        # 清空舊內容後重寫（反映目前狀態，已解決的會自動消失）
        rng = urllib.parse.quote(f"{tab}!A1:D5000")
        _sheet_req("POST", f"{FAQ_SHEET_ID}/values/{rng}:clear", {})
        ts = now_tw_date()
        rows = [header] + [[p["name"], p["reason"], p["size"], ts] for p in _PROBLEMS]
        _sheet_req("PUT",
                   f"{FAQ_SHEET_ID}/values/{urllib.parse.quote(tab + '!A1')}?valueInputOption=RAW",
                   {"values": rows})
        print(f"  📋 索引問題清單已更新：{len(_PROBLEMS)} 個檔案需處理")
    except Exception as e:
        print(f"  ⚠ 寫入索引問題清單失敗：{e}")

# ── 步驟三：索引正式分頁的 FAQ ───────────────────────────────────────────────
def index_faq(index, business):
    ns = business["namespace"]
    # Pinecone 向量 ID 只能是 ASCII：namespace 含中文時用雜湊；ns="" 保留原規則（採購不重複）
    ns_tag = "" if not ns else hashlib.md5(ns.encode()).hexdigest()[:8]
    rows = sheet_read(business["sheet_tab"])
    items = []
    for i in range(1, len(rows)):
        r = (rows[i] + [""] * 10)[:10]
        num, q, a = r[0].strip(), r[3].strip(), r[4].strip()
        if not q or not a:
            continue
        items.append({
            "id": f"faq-{ns_tag}-" + hashlib.md5((num or q).encode()).hexdigest(),
            "text": f"問：{q}\n答：{a}", "source": f"FAQ：{q}",
            "meta_extra": {"faq_answer": a[:3000], "faq_question": q[:500]},
        })
    if items:
        for j in range(0, len(items), EMBED_BATCH):
            bc = items[j:j+EMBED_BATCH]
            upsert(index, ns, bc, embed([c["text"] for c in bc]), "faq")
        print(f"  FAQ 索引 {len(items)} 筆（namespace={ns}）")
    return len(items)

# ── 大綱索引（教學/大型檔：自動生大綱並建索引，提升綜覽問題召回）──────────────
OUTLINE_CKPT = os.path.join(ROOT, "data", "outline_ckpt.json")

def _gemini_outline(text):
    """多模型備援呼叫 Gemini 產生文件大綱（抗 503）。"""
    prompt = ("為以下文件寫一份『大綱索引』，供語意檢索用：條列涵蓋所有主要主題、"
              "流程階段、關鍵名詞與同義說法，讓人用各種問法都能對應到本文件。"
              "300-500字，繁體中文，只輸出大綱。\n\n文件內容：\n" + text[:18000])
    body = json.dumps({"contents": [{"parts": [{"text": prompt}]}],
                       "generationConfig": {"temperature": 0.2, "maxOutputTokens": 1024,
                                            "thinkingConfig": {"thinkingBudget": 0}}}).encode()
    for attempt in range(9):
        model = GEMINI_MODELS[attempt % len(GEMINI_MODELS)]
        try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={GEMINI_API_KEY}"
            req = urllib.request.Request(url, data=body, headers={"Content-Type": "application/json"})
            with urllib.request.urlopen(req, timeout=40) as r:
                j = json.loads(r.read())
            return "".join(p.get("text", "") for p in j["candidates"][0]["content"]["parts"]).strip()
        except Exception:
            time.sleep(8)
    return ""

def reconstruct_from_chunks(index, ns, file_id):
    res = index.query(vector=[0.001] * EMBED_DIM, top_k=1000, include_metadata=True,
                      namespace=ns, filter={"file_id": {"$eq": file_id}})
    chunks = sorted(res.get("matches", []),
                    key=lambda m: m.get("metadata", {}).get("chunk_idx", 0))
    if not chunks:
        return "", ""
    src = chunks[0].get("metadata", {}).get("source", "文件")
    parts = [chunks[0].get("metadata", {}).get("text", "")]
    for m in chunks[1:]:
        parts.append(m.get("metadata", {}).get("text", "")[CHUNK_OVERLAP:])
    return "".join(parts), src

def index_outlines(index, business, drive):
    """為大型教學檔生成大綱並建索引（source_type=outline）。可中斷續跑。"""
    if not GEMINI_API_KEY:
        print("  （無 GEMINI_API_KEY，略過大綱索引）")
        return 0
    ns = business.get("namespace", "")
    try:
        with open(OUTLINE_CKPT) as f:
            done_out = json.load(f)
    except Exception:
        done_out = {}
    files = list_files(drive, business.get("drive_folder_id", ""))
    n = 0
    for f in files:
        fid = f["id"]
        if done_out.get(fid):
            continue
        text, src = reconstruct_from_chunks(index, ns, fid)
        if len(text) < OUTLINE_MIN_CHARS:   # 只處理大型/教學檔
            done_out[fid] = "skip"
            continue
        outline = _gemini_outline(text)
        if outline:
            item = {"id": f"outline-{fid}", "text": outline, "source": src,
                    "meta_extra": {"file_id": fid}}
            upsert(index, ns, [item], embed([outline]), "outline")
            n += 1
            print(f"  📑 大綱索引：{src}")
        done_out[fid] = "done" if outline else "fail"
        with open(OUTLINE_CKPT, "w") as fo:
            json.dump(done_out, fo)
    print(f"  大綱索引完成：新增 {n} 份")
    return n

# ── Main ─────────────────────────────────────────────────────────────────────
def main():
    pc = Pinecone(api_key=PINECONE_API_KEY)
    index = get_index(pc)
    drive = drive_service()
    done = load_ckpt()

    faq_only = os.environ.get("FAQ_ONLY", "") == "1"
    outline_only = os.environ.get("OUTLINE_ONLY", "") == "1"

    if outline_only:
        for b in BUSINESSES:
            print(f"\n=== 大綱索引：{b['name']} ===")
            index_outlines(index, b, drive)
        stats = index.describe_index_stats()
        print(f"\n🎉 大綱索引完成！總向量數：{stats.total_vector_count}")
        return

    for b in BUSINESSES:
        print(f"\n=== 業務：{b['name']}（namespace={b['namespace']}）===")
        # 1) 先同步已審查 FAQ → 正式分頁
        try:
            promote_reviewed(b)
        except Exception as e:
            print(f"  ⚠ 同步已審查失敗：{e}")
        # 2) 先索引 FAQ（快、最高優先，不必等文件）
        try:
            n_faq = index_faq(index, b)
        except Exception as e:
            n_faq = 0
            print(f"  ⚠ FAQ 索引失敗：{e}")
        # 3) FAQ_ONLY 模式：略過耗時的文件索引
        if faq_only:
            print(f"  ✅ {b['name']}（FAQ_ONLY）：FAQ {n_faq} 筆")
            continue
        n_doc = index_drive(index, b, drive, done)
        print(f"  ✅ {b['name']}：FAQ {n_faq} 筆，文件片段 +{n_doc}")

    write_problem_report()   # 把無法索引的檔案列到試算表「索引問題」分頁

    stats = index.describe_index_stats()
    print(f"\n🎉 完成！總向量數：{stats.total_vector_count}")
    print(f"   namespaces：{list((stats.get('namespaces') or {}).keys())}")
    print(f"   時間：{datetime.now(timezone.utc).isoformat()}")

if __name__ == "__main__":
    main()
